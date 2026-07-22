#!/data/data/com.termux/files/usr/bin/python3

import sys
import os
import subprocess
import importlib.util
import time
import curses
import traceback
import zipfile
import tarfile
import csv
import json
import gzip
import shutil
from contextlib import redirect_stderr, redirect_stdout

# --- 1. SETUP & CONFIGURATION ---

# System packages required for Termux (Binaries + Build tools for Python libs)
TERMUX_PACKAGES = [
    "ffmpeg", 
    "unrar", 
    "libcairo", 
    "libgirepository", 
    "libjpeg-turbo", 
    "libpng", 
    "python", 
    "clang",       # Required to build some python libs
    "make",        # Required to build some python libs
    "binutils",    # Required to build some python libs
    "libffi"
]

# (14) Python libraries to auto-install
REQUIRED_MODULES = {
    "Pillow": "Pillow",         # Images
    "reportlab": "reportlab",   # PDF creation
    "docx": "python-docx",    # Word docs
    "odf": "odfpy",           # OpenOffice docs
    "bs4": "beautifulsoup4",  # HTML/XML parsing
    "markdown": "Markdown",     # Markdown parsing
    "lxml": "lxml",           # XML/HTML parsing
    "cairosvg": "cairosvg",     # SVG conversion
    "psd_tools": "psd-tools",   # PSD reading
    "striprtf": "striprtf",     # RTF reading
    "EbookLib": "EbookLib",     # EPUB reading
    "pptx": "python-pptx",    # PowerPoint reading
    "rarfile": "rarfile",       # RAR extraction
    "py7zr": "py7zr"          # 7-Zip extraction
}

# (40) Folders to create
FOLDER_NAMES = [
    "JPG", "PNG", "WEBP", "BMP", "TIFF", "GIF", "ICO", "TGA", "SVG", "PSD",
    "PDF", "TXT", "DOCX", "ODT", "HTML", "MD", "CSV", "RTF", "EPUB", "JSON", "XML", "PPTX",
    "ZIP", "TAR", "RAR", "7Z", "GZ",
    "MP3", "WAV", "OGG", "FLAC", "M4A", "AAC", "WMA",
    "MP4", "MKV", "AVI", "MOV", "WMV", "FLV"
]

IMAGE_FOLDERS = ["JPG", "PNG", "WEBP", "BMP", "TIFF", "GIF", "ICO", "TGA", "SVG", "PSD"]
IMAGE_EXTS = ['.jpg', '.jpeg', '.png', '.webp', '.bmp', '.tiff', '.gif', '.ico', '.tga']
VECTOR_IMAGE_EXTS = ['.svg']
LAYERED_IMAGE_EXTS = ['.psd']
AV_EXTS = ['.mp3', '.wav', '.ogg', '.flac', '.m4a', '.aac', '.wma', '.mp4', '.mkv', '.avi', '.mov', '.wmv', '.flv']
ARCHIVE_EXTS = ['.zip', '.tar', '.gz', '.bz2', '.rar', '.7z']
TEXT_DOC_EXTS = ['.txt', '.docx', '.odt', '.html', '.md', '.csv', '.rtf', '.epub', '.json', '.xml', '.pptx', '.svg']

STORAGE_PATH = "/storage/emulated/0"
DOWNLOAD_PATH = os.path.join(STORAGE_PATH, "Download")
BASE_CONVERTER_PATH = os.path.join(DOWNLOAD_PATH, "File Converter")

HAS_FFMPEG = False
HAS_UNRAR = False
HAS_CAIRO = False

# --- 2. AUTO-INSTALLATION FUNCTIONS ---

def clear_screen_standard():
    os.system('clear')

def install_termux_system_deps():
    """Automatically installs required system packages via pkg."""
    print("--- 1/4 Checking System Dependencies (Termux) ---")
    
    # Check if we are likely in Termux
    if not os.path.exists("/data/data/com.termux/files/usr/bin/pkg"):
        print("Not running in standard Termux environment. Skipping system pkg install.")
        return

    print("Updating package lists and installing binaries...")
    print("This may take a minute. Please wait...")
    
    try:
        # Update repositories
        subprocess.run(["pkg", "update", "-y"], check=False)
        
        # Install packages
        cmd = ["pkg", "install", "-y"] + TERMUX_PACKAGES
        subprocess.run(cmd, check=True)
        print("System binaries installed successfully.\n")
    except Exception as e:
        print(f"WARNING: System install failed: {e}")
        print("Attempting to continue, but some features (FFmpeg, Cairo) may fail.\n")
    time.sleep(1)

def check_and_install_python_deps():
    """Checks and installs required Python modules."""
    print("--- 2/4 Checking Python Libraries ---")
    
    # Upgrade pip first to ensure we can install wheels correctly
    try:
        subprocess.run([sys.executable, "-m", "pip", "install", "--upgrade", "pip"], 
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except:
        pass

    all_installed = True
    for module_name, package_name in REQUIRED_MODULES.items():
        spec = importlib.util.find_spec(module_name)
        if spec is None:
            all_installed = False
            print(f"Installing '{package_name}'... (Building wheels may take time)")
            try:
                # We use subprocess directly to see output if it fails, 
                # but hide it if it succeeds to keep UI clean.
                subprocess.run([sys.executable, "-m", "pip", "install", package_name], check=True)
                print(f"✔ Installed '{package_name}'.")
            except Exception:
                print(f"✖ ERROR: Failed to install '{package_name}'.")
                print(f"  Try manually: pip install {package_name}")
                # We don't exit here, we try to keep going
        else:
            pass
    
    if all_installed:
        print("All Python libraries are present.\n")
    else:
        print("Library check complete.\n")
    time.sleep(0.5)

def check_external_bins_status():
    """Checks which external binaries are actually available after install."""
    global HAS_FFMPEG, HAS_UNRAR, HAS_CAIRO
    print("--- 3/4 Verifying Integrations ---")
    
    # Check ffmpeg
    if shutil.which("ffmpeg"):
        print("✔ 'ffmpeg' detected. A/V conversions ENABLED.")
        HAS_FFMPEG = True
    else:
        print("✖ 'ffmpeg' NOT found. A/V conversions DISABLED.")
        HAS_FFMPEG = False

    # Check unrar
    if shutil.which("unrar"):
        print("✔ 'unrar' detected. RAR extraction ENABLED.")
        HAS_UNRAR = True
    else:
        print("✖ 'unrar' NOT found. RAR extraction DISABLED.")
        HAS_UNRAR = False
        
    # Check cairo
    if importlib.util.find_spec("cairosvg") is not None:
        # Basic check if the library loads without DLL errors
        try:
            import cairosvg
            HAS_CAIRO = True
            print("✔ 'cairosvg' loaded. SVG conversions ENABLED.")
        except OSError:
            print("✖ 'cairosvg' installed but system libs missing. SVG DISABLED.")
            HAS_CAIRO = False
    else:
        print("✖ 'cairosvg' not found. SVG DISABLED.")
        HAS_CAIRO = False
        
    print("")
    time.sleep(0.5)

def ensure_storage_access():
    print("--- 4/4 Checking Storage Access ---")
    if not os.path.exists(DOWNLOAD_PATH):
        print(f"Access to '{DOWNLOAD_PATH}' denied or missing.")
        print("Attempting to request storage permissions...")
        try:
            subprocess.run(["termux-setup-storage"], check=True)
            print("Permission requested. Please allow it in the popup.")
            print("Waiting 5 seconds for permission to propagate...")
            time.sleep(5)
        except FileNotFoundError:
            print("Could not run 'termux-setup-storage'.")
        
        if not os.path.exists(DOWNLOAD_PATH):
             print(f"ERROR: Still cannot access '{DOWNLOAD_PATH}'.")
             print("Please restart Termux and run this script again.")
             sys.exit(1)
             
    print("Storage access confirmed.\n")
    time.sleep(0.5)

def setup_folders():
    try:
        os.makedirs(BASE_CONVERTER_PATH, exist_ok=True)
        for folder in FOLDER_NAMES:
            os.makedirs(os.path.join(BASE_CONVERTER_PATH, folder), exist_ok=True)
    except Exception as e:
        print(f"ERROR: Could not create folders: {e}")
        sys.exit(1)

# --- 3. IMPORTS (Post-Installation) ---
# These are imported inside functions or wrapped to prevent crash before install finishes

# --- 4. CORE CONVERSION LOGIC ---

def get_text_from_file(input_path, in_ext):
    text_lines = []
    try:
        if in_ext == '.txt':
            with open(input_path, 'r', encoding='utf-8') as f:
                text_lines = f.readlines()
        elif in_ext == '.docx':
            from docx import Document
            doc = Document(input_path)
            text_lines = [para.text + '\n' for para in doc.paragraphs]
        elif in_ext == '.odt':
            from odf.opendocument import load as odf_load
            from odf.text import P as odf_P
            doc = odf_load(input_path)
            for para in doc.getElementsByType(odf_P):
                text_lines.append(str(para) + '\n')
        elif in_ext in ['.html', '.xml', '.svg']:
            from bs4 import BeautifulSoup
            with open(input_path, 'r', encoding='utf-8') as f:
                parser = 'lxml' if in_ext != '.html' else 'html.parser'
                soup = BeautifulSoup(f, parser)
                text_lines = [line + '\n' for line in soup.get_text().splitlines()]
        elif in_ext == '.md':
            import markdown
            from bs4 import BeautifulSoup
            with open(input_path, 'r', encoding='utf-8') as f:
                html = markdown.markdown(f.read())
                soup = BeautifulSoup(html, 'html.parser')
                text_lines = [line + '\n' for line in soup.get_text().splitlines()]
        elif in_ext == '.csv':
            with open(input_path, 'r', encoding='utf-8', newline='') as f:
                reader = csv.reader(f)
                text_lines = [','.join(row) + '\n' for row in reader]
        elif in_ext == '.json':
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                text_lines = [json.dumps(data, indent=2)]
        elif in_ext == '.rtf':
            from striprtf.striprtf import rtf_to_text
            with open(input_path, 'r', encoding='utf-8') as f:
                content = f.read()
            text_lines = [rtf_to_text(content)]
        elif in_ext == '.epub':
            from ebooklib import epub, ITEM_DOCUMENT
            from bs4 import BeautifulSoup
            book = epub.read_epub(input_path)
            for item in book.get_items():
                if item.get_type() == ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text_lines.append(soup.get_text() + '\n\n')
        elif in_ext == '.pptx':
            import pptx
            prs = pptx.Presentation(input_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_lines.append(shape.text + '\n')
    except Exception as e:
        raise Exception(f"Text extraction failed: {e}")
    return text_lines

def write_text_to_pdf(text_lines, output_path):
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import inch
    
    c = canvas.Canvas(output_path, pagesize=A4)
    width, height = A4
    margin_x, margin_y = 0.75 * inch, 1 * inch
    text_object = c.beginText(margin_x, height - margin_y)
    text_object.setFont("Helvetica", 10)
    line_height, y = 12, height - margin_y
    for line in text_lines:
        for sub_line in line.split('\n'):
            if y < margin_y:
                c.drawText(text_object)
                c.showPage()
                text_object = c.beginText(margin_x, height - margin_y)
                text_object.setFont("Helvetica", 10)
                y = height - margin_y
            text_object.textLine(sub_line.strip('\r'))
            y -= line_height
    c.drawText(text_object)
    c.save()

def handle_image_conversion(stdscr, in_path, out_path):
    from PIL import Image
    with Image.open(in_path) as img:
        if out_path.lower().endswith(('.jpg', '.jpeg')):
            if img.mode == 'RGBA':
                img = img.convert('RGB')
        img.save(out_path)

def handle_svg_conversion(stdscr, in_path, out_path):
    if not HAS_CAIRO:
        raise Exception("Cairo/SVG libraries not installed/loaded.")
    import cairosvg
    out_ext = os.path.splitext(out_path)[1].lower()
    if out_ext == '.png':
        cairosvg.svg2png(url=in_path, write_to=out_path)
    elif out_ext == '.pdf':
        cairosvg.svg2pdf(url=in_path, write_to=out_path)
    else:
        raise Exception(f"SVG conversion to {out_ext} not supported.")

def handle_psd_conversion(stdscr, in_path, out_path):
    from psd_tools import PSDImage
    psd = PSDImage.open(in_path)
    composite_image = psd.composite()
    composite_image.save(out_path)

def handle_av_conversion(stdscr, in_path, out_path):
    if not HAS_FFMPEG:
        raise Exception("'ffmpeg' not found. A/V conversion is disabled.")
    command = ['ffmpeg', '-i', in_path, '-y', out_path]
    curses.endwin()
    print("--- Running ffmpeg ---")
    print(f"Command: {' '.join(command)}")
    try:
        # Using subprocess directly to show output
        subprocess.run(command, check=True)
        print("ffmpeg conversion finished successfully.")
    except Exception as e:
        print(f"ffmpeg ERROR: {e}")
        raise Exception(f"ffmpeg conversion failed. {e}")
    finally:
        print("Press Enter to return to the app...")
        sys.stdin.read(1)
        stdscr.refresh()

def handle_extraction(stdscr, in_path, out_folder_path, in_ext):
    base_name = os.path.splitext(os.path.basename(in_path))[0]
    extract_path = os.path.join(out_folder_path, base_name)
    os.makedirs(extract_path, exist_ok=True)
    
    if in_ext == '.zip':
        with zipfile.ZipFile(in_path, 'r') as zf:
            zf.extractall(extract_path)
    elif in_ext in ['.tar', '.gz', '.bz2']:
        if in_ext == '.gz' and not in_path.endswith('.tar.gz'):
             out_filename = os.path.splitext(os.path.basename(in_path))[0]
             out_path = os.path.join(out_folder_path, out_filename)
             with gzip.open(in_path, 'rb') as f_in:
                 with open(out_path, 'wb') as f_out:
                     shutil.copyfileobj(f_in, f_out)
             return f"Decompressed to: {out_path}"
        else:
            with tarfile.open(in_path, 'r:*') as tf:
                tf.extractall(extract_path)
    elif in_ext == '.rar':
        if not HAS_UNRAR:
            raise Exception("'unrar' binary not found.")
        import rarfile
        with rarfile.RarFile(in_path) as rf:
            rf.extractall(extract_path)
    elif in_ext == '.7z':
        import py7zr
        with py7zr.SevenZipFile(in_path, 'r') as zf:
            zf.extractall(extract_path)
            
    return f"Extracted to: {extract_path}"

def handle_data_conversion(stdscr, in_path, out_path, in_ext, out_ext):
    if in_ext == '.csv' and out_ext == '.json':
        data = []
        with open(in_path, 'r', encoding='utf-8', newline='') as f:
            reader = csv.DictReader(f)
            for row in reader:
                data.append(row)
        with open(out_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=4)
    elif in_ext == '.json' and out_ext == '.csv':
        with open(in_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        if not isinstance(data, list) or not data:
            raise Exception("JSON must be a non-empty list of objects.")
        headers = data[0].keys()
        with open(out_path, 'w', encoding='utf-8', newline='') as f:
            writer = csv.DictWriter(f, fieldnames=headers)
            writer.writeheader()
            writer.writerows(data)
    else:
        raise Exception(f"Data conversion {in_ext} to {out_ext} not supported.")

def handle_md_to_html(stdscr, in_path, out_path):
    import markdown
    with open(in_path, 'r', encoding='utf-8') as f:
        html = markdown.markdown(f.read())
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write(html)

def handle_single_image_to_pdf(stdscr, in_path, out_path):
    from PIL import Image
    try:
        with Image.open(in_path) as img:
            img_rgb = img.convert('RGB')
            img_rgb.save(out_path, "PDF", resolution=100.0)
    except Exception as e:
        raise Exception(f"Pillow (Image->PDF) Error: {e}")

def handle_multi_image_to_pdf(stdscr, image_paths, out_path):
    from PIL import Image
    try:
        images_rgb = []
        for path in image_paths:
            img = Image.open(path)
            images_rgb.append(img.convert('RGB'))
        if not images_rgb:
            raise Exception("No images found to convert.")
        images_rgb[0].save(
            out_path, "PDF", resolution=100.0,
            save_all=True, append_images=images_rgb[1:]
        )
    except Exception as e:
        raise Exception(f"Pillow (Multi-Image->PDF) Error: {e}")

# --- 5. MAIN CONVERSION ROUTER ---

def convert_file(stdscr, in_path, out_folder_name):
    in_ext = os.path.splitext(in_path)[1].lower()
    out_ext = f".{out_folder_name.lower()}"
    base_name = os.path.splitext(os.path.basename(in_path))[0]
    out_folder_path = os.path.join(BASE_CONVERTER_PATH, out_folder_name)
    out_path = os.path.join(out_folder_path, f"{base_name}{out_ext}")

    try:
        if in_ext in ARCHIVE_EXTS:
            out_folder = out_folder_path if in_ext == '.gz' else os.path.join(BASE_CONVERTER_PATH, out_folder_name)
            message = handle_extraction(stdscr, in_path, out_folder, in_ext)
            return (True, message)

        if in_ext in VECTOR_IMAGE_EXTS and out_ext in ['.png', '.pdf']:
            handle_svg_conversion(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")

        if in_ext in LAYERED_IMAGE_EXTS and out_ext in IMAGE_EXTS:
            handle_psd_conversion(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")

        if in_ext in IMAGE_EXTS and out_ext in IMAGE_EXTS:
            handle_image_conversion(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")
            
        if in_ext in IMAGE_EXTS and out_ext == '.pdf':
            handle_single_image_to_pdf(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")

        if in_ext in AV_EXTS and out_ext in AV_EXTS:
            handle_av_conversion(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")
            
        if in_ext in ['.csv', '.json'] and out_ext in ['.csv', '.json']:
            handle_data_conversion(stdscr, in_path, out_path, in_ext, out_ext)
            return (True, f"Saved to: {out_path}")

        if in_ext == '.md' and out_ext == '.html':
            handle_md_to_html(stdscr, in_path, out_path)
            return (True, f"Saved to: {out_path}")

        if out_ext == '.txt' and in_ext in TEXT_DOC_EXTS:
            text_lines = get_text_from_file(in_path, in_ext)
            with open(out_path, 'w', encoding='utf-8') as f:
                f.writelines(text_lines)
            return (True, f"Saved to: {out_path}")
            
        if out_ext == '.pdf' and in_ext in TEXT_DOC_EXTS:
            text_lines = get_text_from_file(in_path, in_ext)
            write_text_to_pdf(text_lines, out_path)
            return (True, f"Saved to: {out_path}")

        return (False, f"Unsupported conversion: {in_ext} to {out_ext}")

    except Exception as e:
        return (False, f"ERROR: {str(e)}")


# --- 6. CURSES UI HELPER FUNCTIONS ---

def draw_header(stdscr, title):
    h, w = stdscr.getmaxyx()
    stdscr.attron(curses.color_pair(2))
    stdscr.addstr(0, 0, " " * w, curses.color_pair(2))
    header_text = f" Termux Curses Converter (q to quit/back) "
    stdscr.addstr(0, (w - len(header_text)) // 2, header_text, curses.A_REVERSE)
    stdscr.attroff(curses.color_pair(2))
    stdscr.attron(curses.color_pair(1))
    stdscr.addstr(2, (w - len(title)) // 2, title)
    stdscr.attroff(curses.color_pair(1))

def draw_status(stdscr, message, is_error=False):
    h, w = stdscr.getmaxyx()
    nav_hint = " (Arrows/Enter to Select, q to Back) "
    hint_len = len(nav_hint)
    color = curses.color_pair(3) if is_error else curses.color_pair(2)
    stdscr.attron(color)
    stdscr.addstr(h - 1, 0, " " * (w - 1))
    stdscr.addstr(h - 1, 1, message[:w - hint_len - 2])
    stdscr.attron(curses.A_REVERSE)
    stdscr.addstr(h - 1, w - hint_len - 1, nav_hint)
    stdscr.attroff(curses.A_REVERSE)
    stdscr.attroff(color)

def run_menu(stdscr, title, options, sub_title=""):
    current_idx = 0
    while True:
        stdscr.clear()
        h, w = stdscr.getmaxyx()
        draw_header(stdscr, title)
        if sub_title:
            stdscr.addstr(4, (w - len(sub_title)) // 2, sub_title)
        draw_status(stdscr, "Select an option.")
        
        for i, option in enumerate(options):
            display_option = option
            if len(option) > w - 8:
                display_option = option[:w - 11] + "..."
            x = (w - len(display_option)) // 2 - 2
            y = h // 2 - len(options) // 2 + i
            if i == current_idx:
                stdscr.attron(curses.color_pair(1))
                stdscr.addstr(y, x, f"> {display_option} <")
                stdscr.attroff(curses.color_pair(1))
            else:
                stdscr.addstr(y, x, f"  {display_option}  ")
        
        stdscr.refresh()
        key = stdscr.getch()
        if key == curses.KEY_UP:
            current_idx = (current_idx - 1) % len(options)
        elif key == curses.KEY_DOWN:
            current_idx = (current_idx + 1) % len(options)
        elif key in [curses.KEY_ENTER, 10, 13]:
            return options[current_idx]
        elif key == ord('q'):
            return None

def run_file_selector(stdscr, folder_path, title, input_folder_name):
    try:
        files = [f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f))]
        files.sort()
    except Exception as e:
        draw_status(stdscr, f"Error reading {folder_path}: {e}", is_error=True)
        stdscr.getch()
        return None
    
    if not files:
        draw_status(stdscr, f"No files found in {os.path.basename(folder_path)}", is_error=True)
        stdscr.getch()
        return None
        
    options = ["[ .. Go Back .. ]"]
    if input_folder_name in IMAGE_FOLDERS:
        options.append(f"[ ** Convert ALL {len(files)} Images in '{input_folder_name}' to one PDF ** ]")
    options.extend(files)
    selection = run_menu(stdscr, title, options, f"Folder: /Download/File Converter/{input_folder_name}")
    if selection == "[ .. Go Back .. ]":
        return None
    return selection

def run_confirmation(stdscr, prompt):
    options = ["Yes", "No"]
    selection = run_menu(stdscr, prompt, options, "Please confirm")
    return selection

def run_help(stdscr):
    stdscr.clear()
    h, w = stdscr.getmaxyx()
    draw_header(stdscr, "How to Use")
    help_text = [
        "1. MOVE FILES:",
        f"   Go to: /Download/File Converter/",
        "   Move files into the correct INPUT folder.",
        "",
        "2. RUN CONVERTER:",
        "   Select 'Convert a File' -> INPUT folder -> File.",
        "",
        "3. SELECT OUTPUT:",
        "   Select the format you want.",
        "",
        "** NOTES **",
        " - Archives extract to their own folder.",
        " - A/V conversion uses FFmpeg (Auto-installed).",
        " - Images can be combined into one PDF."
    ]
    for i, line in enumerate(help_text):
        if 5 + i >= h - 2: break
        stdscr.addstr(5 + i, (w - len(line)) // 2, line)
    draw_status(stdscr, "Press 'q' or Enter to go back.")
    stdscr.refresh()
    while True:
        key = stdscr.getch()
        if key in [ord('q'), curses.KEY_ENTER, 10, 13]:
            return

def run_text_input(stdscr, prompt):
    stdscr.clear()
    h, w = stdscr.getmaxyx()
    draw_header(stdscr, "Input Required")
    stdscr.addstr(h // 2 - 1, (w - len(prompt)) // 2, prompt)
    box_y, box_x = h // 2 + 1, w // 2 - 20
    stdscr.attron(curses.color_pair(2))
    stdscr.addstr(box_y, box_x, " " * 40)
    stdscr.attroff(curses.color_pair(2))
    draw_status(stdscr, "Type filename. Enter to Confirm.")
    curses.curs_set(1)
    stdscr.keypad(True)
    text = ""
    while True:
        stdscr.attron(curses.color_pair(2))
        stdscr.addstr(box_y, box_x, " " * 40)
        stdscr.addstr(box_y, box_x + 1, text[:38])
        stdscr.attroff(curses.color_pair(2))
        stdscr.move(box_y, box_x + 1 + len(text))
        stdscr.refresh()
        key = stdscr.getch()
        if key in [curses.KEY_ENTER, 10, 13]:
            break
        elif key == ord('q'):
            text = None; break
        elif key in [curses.KEY_BACKSPACE, 127, 8]:
            text = text[:-1]
        elif 32 <= key <= 126:
            if len(text) < 38: text += chr(key)
    curses.curs_set(0); stdscr.keypad(False); return text

def run_multi_image_to_pdf_wizard(stdscr, input_folder_path, input_folder_name):
    try:
        image_paths = [
            os.path.join(input_folder_path, f) 
            for f in os.listdir(input_folder_path) 
            if os.path.splitext(f)[1].lower() in IMAGE_EXTS
        ]
        image_paths.sort()
    except Exception as e:
        draw_status(stdscr, f"Error: {e}", is_error=True); stdscr.getch(); return
    if not image_paths:
        draw_status(stdscr, "No images found.", is_error=True); stdscr.getch(); return
    confirm = run_confirmation(stdscr, f"Combine {len(image_paths)} images into one PDF?")
    if confirm != "Yes": return
    default_name = f"{input_folder_name}_Album"
    filename = run_text_input(stdscr, f"Enter PDF name (default: {default_name})")
    if filename is None: return
    if not filename: filename = default_name
    out_folder_path = os.path.join(BASE_CONVERTER_PATH, "PDF")
    out_path = os.path.join(out_folder_path, f"{filename}.pdf")
    draw_status(stdscr, "Working..."); stdscr.refresh()
    try:
        handle_multi_image_to_pdf(stdscr, image_paths, out_path)
        draw_status(stdscr, f"Success! Saved to: /PDF/{filename}.pdf")
    except Exception as e:
        draw_status(stdscr, f"ERROR: {e}", is_error=True)
    stdscr.getch()

# --- 7. MAIN CURSES APPLICATION ---

def main(stdscr):
    curses.curs_set(0); stdscr.nodelay(0); stdscr.timeout(-1)
    curses.start_color()
    curses.init_pair(1, curses.COLOR_CYAN, curses.COLOR_BLACK)
    curses.init_pair(2, curses.COLOR_BLACK, curses.COLOR_WHITE)
    curses.init_pair(3, curses.COLOR_WHITE, curses.COLOR_RED)

    while True:
        main_choice = run_menu(stdscr, "Main Menu", ["Convert a File", "Help / How to Use", "Exit"])
        if main_choice == "Exit" or main_choice is None: break
        if main_choice == "Help / How to Use": run_help(stdscr); continue
            
        input_folder = run_menu(stdscr, "Step 1: Choose INPUT Folder", FOLDER_NAMES)
        if input_folder is None: continue
        input_folder_path = os.path.join(BASE_CONVERTER_PATH, input_folder)
        input_file = run_file_selector(stdscr, input_folder_path, f"Step 2: Choose File from '{input_folder}'", input_folder)
        if input_file is None: continue
        
        if input_file.startswith("[ ** Convert ALL"):
            run_multi_image_to_pdf_wizard(stdscr, input_folder_path, input_folder)
            continue
            
        full_input_path = os.path.join(input_folder_path, input_file)
        in_ext = os.path.splitext(input_file)[1].lower()
        if in_ext in ARCHIVE_EXTS:
            output_folder = input_folder
            prompt = f"Extract '{input_file}' here?"
        else:
            output_folder = run_menu(stdscr, "Step 3: Choose OUTPUT Format", FOLDER_NAMES)
            if output_folder is None: continue
            if output_folder == input_folder:
                draw_status(stdscr, "Input and Output cannot be the same.", is_error=True)
                stdscr.getch(); continue
            prompt = f"Convert '{input_file}' to {output_folder}?"
             
        confirm = run_confirmation(stdscr, prompt)
        if confirm != "Yes": continue

        draw_status(stdscr, "Working, please wait..."); stdscr.refresh()
        success, message = convert_file(stdscr, full_input_path, output_folder)
        draw_status(stdscr, message, is_error=not success)
        stdscr.getch()

# --- 8. SCRIPT ENTRYPOINT ---

if __name__ == "__main__":
    try:
        clear_screen_standard()
        print("--- Initializing Termux Converter v3.2 (Auto-Install) ---")
        
        # 1. Install System Binaries (pkg)
        install_termux_system_deps()
        
        # 2. Install Python Modules (pip)
        check_and_install_python_deps()
        
        # 3. Verify Integrations
        check_external_bins_status()
        
        # 4. Check Storage
        ensure_storage_access()
        
        setup_folders()
        
        print("--- Setup Complete ---")
        print(f"Folders ready in: /storage/emulated/0/Download/File Converter/")
        print("\nStarting application...")
        time.sleep(1)
        
        curses.wrapper(main)
        print("File Converter exited successfully.")

    except KeyboardInterrupt:
        print("\nExiting...")
    except Exception as e:
        print("\nA critical error occurred:")
        traceback.print_exc()
    finally:
        os.system('clear')